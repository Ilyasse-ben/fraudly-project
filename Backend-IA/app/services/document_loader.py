"""
Document Loader — extraction de texte depuis PDF, DOCX, PPTX.

Règle d'architecture :
    Ce fichier est appelé UNIQUEMENT par rag_service.py.
    Jamais directement par une API, un agent ou un consumer Kafka.

Stratégies PDF (PDF_IMAGE_STRATEGY dans .env) :
    "text_only"  → texte brut uniquement, sans OCR          (défaut)
    "ocr_images" → texte + OCR sur les images embarquées    (recommandé prod)
    "ocr_full"   → détecte pages scannées et rasterise
    "full"       → ocr_images + ocr_full combinés

Dépendances système :
    apt-get install -y tesseract-ocr tesseract-ocr-fra tesseract-ocr-eng

Dépendances Python :
    pymupdf, python-docx, python-pptx, pytesseract, Pillow
"""

import io
import re
import fitz
import pytesseract
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation
from typing import List, Dict, Any, Optional

from app.schemas.common import DocumentType
from app.core.config import settings
from app.core.logger import get_logger

logger = get_logger(__name__)

# ── Config Tesseract ──────────────────────────────────────
TESSERACT_LANG   = getattr(settings, "OCR_LANG", "fra+eng")
TESSERACT_BLOCK  = "--oem 3 --psm 6"
TESSERACT_SPARSE = "--oem 3 --psm 11"

MIN_IMAGE_WIDTH        = 100
MIN_IMAGE_HEIGHT       = 100
RASTER_ZOOM            = 1.5
MAX_RASTER_WIDTH       = 2000
SCANNED_PAGE_THRESHOLD = 50


# ── Nettoyage texte ───────────────────────────────────────

def _clean_text(text: str) -> str:
    """
    Normalise le texte avant indexation.
    Texte propre → embedding plus précis → meilleur retrieval.
    """
    text = re.sub(r"[^\S\n\t ]+", " ", text)   # caractères de contrôle
    text = re.sub(r"[ \t]{2,}", " ", text)       # espaces multiples
    text = re.sub(r"\n{3,}", "\n\n", text)       # sauts de ligne excessifs
    return text.strip()


# ── Utilitaires OCR ───────────────────────────────────────

def _bytes_to_pil(image_bytes: bytes) -> Image.Image:
    return Image.open(io.BytesIO(image_bytes))


def _ocr_image_bytes(image_bytes: bytes, sparse: bool = False) -> str:
    """
    Lance Tesseract. Retourne "" silencieusement en cas d'erreur.
    Ne jamais bloquer l'ingestion pour une image.
    """
    try:
        pil_img = _bytes_to_pil(image_bytes)
        config  = TESSERACT_SPARSE if sparse else TESSERACT_BLOCK
        text    = pytesseract.image_to_string(pil_img, lang=TESSERACT_LANG, config=config)
        return _clean_text(text)
    except pytesseract.TesseractNotFoundError:
        logger.error(
            "[OCR] Tesseract non trouvé. "
            "Installer : apt-get install tesseract-ocr tesseract-ocr-fra tesseract-ocr-eng"
        )
        return ""
    except Exception as e:
        logger.warning(f"[OCR] Erreur Tesseract : {e}")
        return ""


def _ocr_page_raster(page: fitz.Page) -> str:
    """
    Rasterise une page PDF complète et OCRise.
    alpha=False → RAM -25%. Zoom adaptatif → protège la RAM sur grands formats.
    """
    page_width_pt = page.rect.width
    natural_px    = page_width_pt * RASTER_ZOOM
    zoom          = RASTER_ZOOM if natural_px <= MAX_RASTER_WIDTH else MAX_RASTER_WIDTH / page_width_pt

    pix         = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    image_bytes = pix.tobytes("png")
    return _ocr_image_bytes(image_bytes, sparse=False)


def _extract_images_from_page(
    doc: fitz.Document,
    page: fitz.Page,
    page_num: int,
) -> List[str]:
    """OCRise les images embarquées d'une page PDF."""
    ocr_results = []
    for img_info in page.get_images(full=True):
        xref, _, width, height = img_info[0], img_info[1], img_info[2], img_info[3]
        if width < MIN_IMAGE_WIDTH or height < MIN_IMAGE_HEIGHT:
            continue
        try:
            image_bytes = doc.extract_image(xref)["image"]
            is_figure   = (width > 300 and height > 200)
            text_ocr    = _ocr_image_bytes(image_bytes, sparse=is_figure)
            if text_ocr:
                label = "Graphique" if is_figure else "Image"
                ocr_results.append(f"[{label} page {page_num}, {width}x{height}px]\n{text_ocr}")
        except Exception as e:
            logger.warning(f"[Loader] Image xref={xref} p.{page_num} ignorée : {e}")
    return ocr_results


# ── PDF — 4 stratégies ────────────────────────────────────

def load_pdf(file_path: str) -> List[Dict[str, Any]]:
    """text_only — texte brut, pas d'OCR."""
    pages = []
    with fitz.open(file_path) as doc:
        logger.info(f"[Loader] PDF text_only — {len(doc)} pages")
        for i, page in enumerate(doc):
            text = _clean_text(page.get_text("text", sort=True))
            if text:
                pages.append({"text": text, "page_number": i + 1})
    return pages


def load_pdf_ocr_images(file_path: str) -> List[Dict[str, Any]]:
    """ocr_images — texte + OCR images embarquées."""
    pages = []
    with fitz.open(file_path) as doc:
        logger.info(f"[Loader] PDF ocr_images — {len(doc)} pages")
        for i, page in enumerate(doc):
            page_num    = i + 1
            text        = _clean_text(page.get_text("text", sort=True))
            image_texts = _extract_images_from_page(doc, page, page_num)
            parts       = ([text] if text else []) + image_texts
            full        = _clean_text("\n\n".join(parts))
            if full:
                pages.append({"text": full, "page_number": page_num})
    return pages


def load_pdf_ocr_full(file_path: str) -> List[Dict[str, Any]]:
    """ocr_full — rasterise les pages scannées."""
    pages = []
    with fitz.open(file_path) as doc:
        logger.info(f"[Loader] PDF ocr_full — {len(doc)} pages")
        for i, page in enumerate(doc):
            page_num = i + 1
            text     = _clean_text(page.get_text("text", sort=True))
            if len(text) >= SCANNED_PAGE_THRESHOLD:
                pages.append({"text": text, "page_number": page_num})
            else:
                ocr_text = _ocr_page_raster(page)
                if ocr_text:
                    pages.append({
                        "text":        f"[Page scannée {page_num}]\n{ocr_text}",
                        "page_number": page_num,
                    })
    return pages


def load_pdf_full(file_path: str) -> List[Dict[str, Any]]:
    """full — texte + images + pages scannées. Production."""
    pages = []
    with fitz.open(file_path) as doc:
        logger.info(f"[Loader] PDF full — {len(doc)} pages")
        for i, page in enumerate(doc):
            page_num = i + 1
            text     = _clean_text(page.get_text("text", sort=True))
            parts: List[str] = []
            if len(text) >= SCANNED_PAGE_THRESHOLD:
                parts.append(text)
                parts.extend(_extract_images_from_page(doc, page, page_num))
            else:
                ocr_text = _ocr_page_raster(page)
                if ocr_text:
                    parts.append(f"[Page scannée {page_num}]\n{ocr_text}")
            full = _clean_text("\n\n".join(parts))
            if full:
                pages.append({"text": full, "page_number": page_num})
    return pages


# ── DOCX ─────────────────────────────────────────────────

def _get_docx_image_page(rel_id: str, doc: DocxDocument, block_size: int) -> Optional[int]:
    """
    Tente d'estimer le numéro de page d'une image DOCX en retrouvant
    sa position ordinale parmi les relations de type image.
    Retourne None si non déterminable (les images seront groupées en page dédiée).
    """
    image_rels = [
        rid for rid, rel in doc.part.rels.items()
        if "image" in rel.reltype
    ]
    try:
        idx = image_rels.index(rel_id)
        return (idx // block_size) + 1
    except ValueError:
        return None


def load_docx(file_path: str) -> List[Dict[str, Any]]:
    """
    Paragraphes + OCR images embarquées.
    FIX : les images sont maintenant associées à leur bloc de page estimé
    plutôt que groupées sur une fausse page N+1 unique.
    """
    doc        = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    logger.info(f"[Loader] DOCX — {len(paragraphs)} paragraphes")

    BLOCK = 50

    # Construire les pages texte d'abord
    pages: List[Dict[str, Any]] = [
        {
            "text":        _clean_text("\n".join(paragraphs[i : i + BLOCK])),
            "page_number": (i // BLOCK) + 1,
        }
        for i in range(0, len(paragraphs), BLOCK)
        if paragraphs[i : i + BLOCK]
    ]

    # Préparer un dict page_number → liste de textes OCR à fusionner
    image_texts_by_page: Dict[int, List[str]] = {}

    for rel_id, rel in doc.part.rels.items():
        if "image" not in rel.reltype:
            continue
        try:
            image_bytes = rel.target_part.blob
            pil_img     = _bytes_to_pil(image_bytes)
            w, h        = pil_img.size
            if w < MIN_IMAGE_WIDTH or h < MIN_IMAGE_HEIGHT:
                continue
            ocr_text = _ocr_image_bytes(image_bytes, sparse=(w > 300 and h > 200))
            if not ocr_text:
                continue
            label = "Graphique" if (w > 300 and h > 200) else "Image"

            # FIX : associer l'image à la page estimée, ou créer une page dédiée
            estimated_page = _get_docx_image_page(rel_id, doc, BLOCK)
            target_page    = estimated_page if estimated_page else (len(pages) + 1)

            image_texts_by_page.setdefault(target_page, []).append(
                f"[{label} DOCX {w}x{h}px]\n{ocr_text}"
            )
        except Exception as e:
            logger.warning(f"[Loader] Image DOCX ignorée : {e}")

    # Fusionner les images OCR dans les pages existantes ou créer des pages dédiées
    for page_num, img_texts in image_texts_by_page.items():
        existing = next((p for p in pages if p["page_number"] == page_num), None)
        img_block = _clean_text("\n\n".join(img_texts))
        if existing:
            existing["text"] = _clean_text(existing["text"] + "\n\n" + img_block)
        else:
            pages.append({"text": img_block, "page_number": page_num})

    # Re-trier par numéro de page après fusion
    pages.sort(key=lambda p: p["page_number"])
    return pages


# ── PPTX ─────────────────────────────────────────────────

def load_pptx(file_path: str) -> List[Dict[str, Any]]:
    """Texte slide + OCR images par slide."""
    prs   = Presentation(file_path)
    pages = []
    logger.info(f"[Loader] PPTX — {len(prs.slides)} slides")

    for i, slide in enumerate(prs.slides):
        page_num = i + 1
        parts: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(r.text for r in para.runs).strip()
                    if line:
                        parts.append(line)

        for shape in slide.shapes:
            if shape.shape_type == 13:
                try:
                    image_bytes = shape.image.blob
                    pil_img     = _bytes_to_pil(image_bytes)
                    w, h        = pil_img.size
                    if w >= MIN_IMAGE_WIDTH and h >= MIN_IMAGE_HEIGHT:
                        ocr_text = _ocr_image_bytes(image_bytes, sparse=(w > 300 and h > 200))
                        if ocr_text:
                            parts.append(f"[Image slide {page_num}, {w}x{h}px]\n{ocr_text}")
                except Exception as e:
                    logger.warning(f"[Loader] Image slide {page_num} ignorée : {e}")

        if parts:
            pages.append({
                "text":        _clean_text("\n".join(parts)),
                "page_number": page_num,
            })
    return pages


# ── Dispatch — appelé UNIQUEMENT par rag_service.py ──────

_PDF_STRATEGIES = {
    "text_only":  load_pdf,
    "ocr_images": load_pdf_ocr_images,
    "ocr_full":   load_pdf_ocr_full,
    "full":       load_pdf_full,
}

_VALID_STRATEGIES = set(_PDF_STRATEGIES.keys())


def load_document(file_path: str, doc_type: DocumentType) -> List[Dict[str, Any]]:
    """
    Point d'entrée unique.
    Retourne toujours : [{"text": str, "page_number": int}, ...]
    FIX : log warning explicite si PDF_IMAGE_STRATEGY est inconnue (faute de frappe .env).
    """
    if doc_type == DocumentType.PDF:
        strategy = getattr(settings, "PDF_IMAGE_STRATEGY", "text_only")
        if strategy not in _VALID_STRATEGIES:
            logger.warning(
                f"[Loader] Stratégie PDF inconnue : '{strategy}'. "
                f"Valeurs acceptées : {sorted(_VALID_STRATEGIES)}. "
                "Fallback sur 'text_only'."
            )
        loader = _PDF_STRATEGIES.get(strategy, load_pdf)
        logger.info(f"[Loader] Stratégie PDF : '{strategy}'")
        return loader(file_path)
    elif doc_type == DocumentType.DOCX:
        return load_docx(file_path)
    elif doc_type == DocumentType.PPTX:
        return load_pptx(file_path)
    else:
        raise ValueError(f"Type de document non supporté : {doc_type}")